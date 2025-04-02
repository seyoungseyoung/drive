import requests
import json
from flask import Flask, request, jsonify, session, render_template, send_file
import logging
import os
from dotenv import load_dotenv
import uuid
import time
from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import io
from PIL import Image
import base64
from werkzeug.utils import secure_filename

# .env 파일 로드
load_dotenv()

# 로깅 설정
logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(name)s - %(levelname)s - %(message)s',
    handlers=[
        logging.FileHandler('server.log'),
        logging.StreamHandler()
    ]
)
logger = logging.getLogger('ppt_agent')

app = Flask(__name__)
app.secret_key = os.getenv('FLASK_SECRET_KEY', 'default_secret_key')
app.config['MAX_CONTENT_LENGTH'] = 16 * 1024 * 1024
app.config['UPLOAD_FOLDER'] = 'uploads'
app.config['ALLOWED_EXTENSIONS'] = {'pptx', 'pdf', 'jpg', 'jpeg', 'png', 'gif', 'svg'}

# 세션 데이터 저장
ppt_sessions = {}

# 이미지 업로드 설정
UPLOAD_FOLDER = 'static/uploads'
ALLOWED_EXTENSIONS = {'png', 'jpg', 'jpeg', 'gif'}

app.config['UPLOAD_FOLDER'] = UPLOAD_FOLDER

def allowed_file(filename):
    return '.' in filename and filename.rsplit('.', 1)[1].lower() in ALLOWED_EXTENSIONS

@app.route('/')
def index():
    if 'session_id' not in session:
        session['session_id'] = str(uuid.uuid4())
        ppt_sessions[session['session_id']] = {
            'slides': [],
            'theme': 'modern',
            'color_palette': 'blue',
            'font_family': 'Pretendard',
            'extensions': {
                'ai_enabled': False,
                'ai_suggestions': [],
                'version_history': []
            }
        }
    return render_template('index.html')

@app.route('/save_slides', methods=['POST'])
def save_slides():
    """슬라이드 저장"""
    try:
        session_id = session.get('session_id')
        if not session_id:
            return jsonify({'success': False, 'error': 'No session ID found'}), 400
            
        data = request.get_json()
        logger.info(f"Saving slides: {json.dumps(data)[:100]}...")
        
        if not data or 'slides' not in data:
            return jsonify({'success': False, 'error': 'No slide data provided'}), 400
        
        # 슬라이드 배열 전체 저장
        ppt_sessions[session_id]['slides'] = data['slides']
        
        # 버전 히스토리 저장 (AI 확장 기능이 활성화된 경우)
        if ppt_sessions[session_id]['extensions']['ai_enabled']:
            ppt_sessions[session_id]['extensions']['version_history'].append({
                'slides': data['slides'],
                'timestamp': time.time()
            })
        
        return jsonify({'success': True})
        
    except Exception as e:
        logger.error(f"Error saving slides: {str(e)}")
        return jsonify({'success': False, 'error': str(e)}), 500

@app.route('/get_slides', methods=['GET'])
def get_slides():
    """모든 슬라이드 데이터 가져오기"""
    try:
        session_id = session.get('session_id')
        if not session_id:
            return jsonify({'success': False, 'error': 'No session ID found'}), 400
            
        # 세션이 존재하지 않으면 생성
        if session_id not in ppt_sessions:
            ppt_sessions[session_id] = {
                'slides': [],
                'theme': 'modern',
                'color_palette': 'blue',
                'font_family': 'Pretendard',
                'extensions': {
                    'ai_enabled': False,
                    'ai_suggestions': [],
                    'version_history': []
                }
            }
        
        # 슬라이드가 없으면 빈 배열 반환
        slides = ppt_sessions[session_id].get('slides', [])
        logger.info(f"Getting slides. Count: {len(slides)}")
        
        return jsonify({
            'success': True,
            'slides': slides,
            'theme': ppt_sessions[session_id].get('theme', 'modern'),
            'color_palette': ppt_sessions[session_id].get('color_palette', 'blue'),
            'font_family': ppt_sessions[session_id].get('font_family', 'Pretendard'),
            'extensions': ppt_sessions[session_id].get('extensions', {})
        })
        
    except Exception as e:
        logger.error(f"Error getting slides: {str(e)}")
        return jsonify({'success': False, 'error': str(e)}), 500

@app.route('/update_theme', methods=['POST'])
def update_theme():
    """프레젠테이션 테마 업데이트"""
    try:
        session_id = session.get('session_id')
        if not session_id:
            return jsonify({'error': 'No session ID found'}), 400
            
        data = request.get_json()
        theme = data.get('theme')
        
        if not theme:
            return jsonify({'error': 'Theme is required'}), 400
        
        ppt_sessions[session_id]['theme'] = theme
        
        return jsonify({'success': True})
        
    except Exception as e:
        logger.error(f"Error updating theme: {str(e)}")
        return jsonify({'error': str(e)}), 500

@app.route('/update_color_palette', methods=['POST'])
def update_color_palette():
    """프레젠테이션 색상 팔레트 업데이트"""
    try:
        session_id = session.get('session_id')
        if not session_id:
            return jsonify({'error': 'No session ID found'}), 400
            
        data = request.get_json()
        color_palette = data.get('color_palette')
        
        if not color_palette:
            return jsonify({'error': 'Color palette is required'}), 400
        
        ppt_sessions[session_id]['color_palette'] = color_palette
        
        return jsonify({'success': True})
        
    except Exception as e:
        logger.error(f"Error updating color palette: {str(e)}")
        return jsonify({'error': str(e)}), 500

@app.route('/update_font', methods=['POST'])
def update_font():
    """프레젠테이션 글꼴 업데이트"""
    try:
        session_id = session.get('session_id')
        if not session_id:
            return jsonify({'error': 'No session ID found'}), 400
            
        data = request.get_json()
        font_family = data.get('font_family')
        
        if not font_family:
            return jsonify({'error': 'Font family is required'}), 400
        
        ppt_sessions[session_id]['font_family'] = font_family
        
        return jsonify({'success': True})
        
    except Exception as e:
        logger.error(f"Error updating font: {str(e)}")
        return jsonify({'error': str(e)}), 500

@app.route('/toggle_extension', methods=['POST'])
def toggle_extension():
    """확장 기능 활성화/비활성화"""
    try:
        session_id = session.get('session_id')
        if not session_id:
            return jsonify({'error': 'No session ID found'}), 400
            
        data = request.get_json()
        extension_name = data.get('extension')
        enabled = data.get('enabled', False)
        
        if extension_name == 'ai':
            ppt_sessions[session_id]['extensions']['ai_enabled'] = enabled
        
        return jsonify({'success': True})
        
    except Exception as e:
        logger.error(f"Error toggling extension: {str(e)}")
        return jsonify({'error': str(e)}), 500

@app.route('/export', methods=['POST'])
def export():
    """프레젠테이션 내보내기"""
    try:
        session_id = session.get('session_id')
        if not session_id:
            return jsonify({'error': 'No session ID found'}), 400
            
        data = request.get_json()
        format = data.get('format', 'pptx')
        
        if format not in ['pptx', 'pdf']:
            return jsonify({'error': 'Unsupported format'}), 400
        
        presentation_data = ppt_sessions.get(session_id)
        if not presentation_data or not presentation_data.get('slides'):
            return jsonify({'error': 'No presentation data found'}), 400
        
        if format == 'pptx':
            # PowerPoint 파일 생성
            prs = Presentation()
            
            # 테마 적용
            theme = presentation_data['theme']
            color_palette = presentation_data['color_palette']
            font_family = presentation_data['font_family']
            
            # 각 슬라이드 생성
            for slide_data in presentation_data['slides']:
                slide = prs.slides.add_slide(prs.slide_layouts[6])
                
                # 배경 설정
                if 'background' in slide_data:
                    background = slide.background
                    if slide_data['background']['type'] == 'solid':
                        background.fill.solid()
                        background.fill.fore_color.rgb = RGBColor.from_string(slide_data['background']['color'])
                
                # 요소 추가
                for element in slide_data.get('elements', []):
                    if element['type'] == 'text':
                        text_box = slide.shapes.add_textbox(
                            element['x'],
                            element['y'],
                            element['width'],
                            element['height']
                        )
                        text_frame = text_box.text_frame
                        p = text_frame.add_paragraph()
                        p.text = element['content']
                        if 'style' in element:
                            if 'fontSize' in element['style']:
                                p.font.size = Pt(element['style']['fontSize'])
                            if 'fontWeight' in element['style']:
                                p.font.bold = element['style']['fontWeight'] == 'bold'
                    elif element['type'] == 'shape':
                        shape = slide.shapes.add_shape(
                            MSO_SHAPE.RECTANGLE,
                            element['x'],
                            element['y'],
                            element['width'],
                            element['height']
                        )
                        if 'style' in element:
                            if 'color' in element['style']:
                                shape.fill.solid()
                                shape.fill.fore_color.rgb = RGBColor.from_string(element['style']['color'])
                            if 'opacity' in element['style']:
                                shape.fill.fore_color.alpha = int(element['style']['opacity'] * 100000)
            
            # 메모리에 저장
            pptx_io = io.BytesIO()
            prs.save(pptx_io)
            pptx_io.seek(0)
            
            return send_file(
                pptx_io,
                mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation',
                as_attachment=True,
                download_name='presentation.pptx'
            )
            
        elif format == 'pdf':
            # PDF 내보내기 구현
            return jsonify({'error': 'PDF export not implemented yet'}), 501
        
    except Exception as e:
        logger.error(f"Error exporting presentation: {str(e)}")
        return jsonify({'error': str(e)}), 500

@app.route('/upload_image', methods=['POST'])
def upload_image():
    if 'image' not in request.files:
        return jsonify({'success': False, 'error': '이미지 파일이 없습니다.'})
    
    file = request.files['image']
    if file.filename == '':
        return jsonify({'success': False, 'error': '선택된 파일이 없습니다.'})
    
    if file and allowed_file(file.filename):
        filename = secure_filename(file.filename)
        filepath = os.path.join(app.config['UPLOAD_FOLDER'], filename)
        file.save(filepath)
        
        return jsonify({
            'success': True,
            'image_url': f'/static/uploads/{filename}'
        })
    
    return jsonify({'success': False, 'error': '허용되지 않는 파일 형식입니다.'})

if __name__ == '__main__':
    # 디렉토리 확인
    os.makedirs('uploads', exist_ok=True)
    os.makedirs('static', exist_ok=True)
    
    # 서버 설정
    host = '0.0.0.0'  # 모든 IP에서 접근 가능
    port = 5000
    debug = True
    
    logger.info(f"Starting server on {host}:{port}")
    app.run(host=host, port=port, debug=debug)

